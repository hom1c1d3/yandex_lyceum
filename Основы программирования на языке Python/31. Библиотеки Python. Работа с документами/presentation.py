import random
from pptx import Presentation
from pptx.util import Pt

prs = Presentation()
title_content_slide_layout = prs.slide_layouts[1]
methods = random.sample([i for i in dir(random) if i.islower() and not i.startswith('_')], 5)

for i in methods:
    slide = prs.slides.add_slide(title_content_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = i
    p = body_shape.text_frame.paragraphs[-1]
    run = p.add_run()
    run.text = eval(f'random.{i}.__doc__')
    run.font.name = 'Courier New'
    run.font.size = Pt(10)


prs.save('test.pptx')